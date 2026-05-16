"""
Reproduce the ProvMind system figure as a PowerPoint slide.
Run:  python3 gen_figure_pptx.py
Output: provmind_figure.pptx  (same directory)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── colour palette ──────────────────────────────────────────────────────────
C_PANEL_BORDER  = RGBColor(0x2B, 0x7A, 0x8C)   # teal border
C_PANEL_BG      = RGBColor(0xF4, 0xF9, 0xFC)   # near-white panel fill
C_HDR_BLUE      = RGBColor(0x1F, 0x5C, 0x8B)   # panel-A header dark blue
C_HDR_TEAL      = RGBColor(0x2B, 0x7A, 0x8C)   # other panel headers
C_HDR_ORANGE    = RGBColor(0xE0, 0x70, 0x20)   # sub-section headers D
C_NODE_GREEN    = RGBColor(0x54, 0x82, 0x35)   # activity nodes (green ellipse)
C_BOX_BLUE_FILL = RGBColor(0xBD, 0xD7, 0xEE)   # light-blue entity boxes
C_BOX_BLUE_BDR  = RGBColor(0x2E, 0x75, 0xB6)   # border of entity boxes
C_DIAMOND_ORG   = RGBColor(0xF4, 0xA2, 0x61)   # orange diamond (tool)
C_DB_PURPLE     = RGBColor(0x70, 0x30, 0xA0)   # database cylinder accent
C_GREEN_TICK    = RGBColor(0x37, 0x86, 0x10)   # correct-answer green
C_BAR_BLUE      = RGBColor(0x2E, 0x75, 0xB6)   # progress bar
C_WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK         = RGBColor(0x00, 0x00, 0x00)
C_LIGHT_GRAY    = RGBColor(0xD9, 0xD9, 0xD9)
C_GOLD          = RGBColor(0xFF, 0xC0, 0x00)
C_SECTION_A     = RGBColor(0x4E, 0xA7, 0xC4)   # subsection row A
C_SECTION_B     = RGBColor(0xED, 0x7D, 0x31)   # subsection row B
C_SECTION_C     = RGBColor(0xFF, 0xC0, 0x00)   # subsection row C
C_SECTION_D2    = RGBColor(0x70, 0xAD, 0x47)   # subsection row D (order)

# ── slide dimensions  (widescreen 20 × 12 inch) ────────────────────────────
W = Inches(20)
H = Inches(12)

# ── helpers ─────────────────────────────────────────────────────────────────

def new_slide(prs):
    layout = prs.slide_layouts[6]          # blank
    return prs.slides.add_slide(layout)


def add_rect(slide, x, y, w, h,
             fill=None, border=None, border_w=Pt(1.5),
             radius=None):
    """Add a plain rectangle (optionally rounded)."""
    shape = slide.shapes.add_shape(
        1,                                 # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.line.width = border_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
    else:
        shape.line.fill.background()
    if radius is not None:
        # set rounded corners via XML
        sp = shape.element
        prstGeom = sp.find(qn('p:spPr')).find(qn('a:prstGeom'))
        if prstGeom is not None:
            prstGeom.set('prst', 'roundRect')
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            else:
                for child in list(avLst):
                    avLst.remove(child)
            gd = etree.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', f'val {radius}')
    return shape


def add_ellipse(slide, x, y, w, h, fill=None, border=None, border_w=Pt(1.2)):
    shape = slide.shapes.add_shape(9, x, y, w, h)   # oval
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = border_w
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, font_size=Pt(9),
                bold=False, color=C_BLACK, align=PP_ALIGN.LEFT,
                wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_label_box(slide, x, y, w, h, letter, title,
                  hdr_color=C_HDR_TEAL, hdr_h=Inches(0.35)):
    """Draw a panel with a coloured header bar containing 'A  Title text'."""
    # outer border
    add_rect(slide, x, y, w, h,
             fill=C_PANEL_BG, border=C_PANEL_BORDER, border_w=Pt(2),
             radius=25000)
    # header bar
    hdr = add_rect(slide, x, y, w, hdr_h,
                   fill=hdr_color, border=None, radius=25000)
    # header text
    add_textbox(slide,
                x + Inches(0.12), y + Inches(0.02),
                w - Inches(0.15), hdr_h,
                f"{letter}   {title}",
                font_size=Pt(11), bold=True,
                color=C_WHITE, align=PP_ALIGN.LEFT)
    return hdr


def add_arrow(slide, x1, y1, x2, y2, color=C_BLACK, width=Pt(1.2)):
    """Add a straight connector arrow from (x1,y1) to (x2,y2)."""
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    # end arrow
    cxnSp = connector.element
    spPr = cxnSp.find(qn('p:spPr'))
    ln = spPr.find(qn('a:ln'))
    if ln is None:
        ln = etree.SubElement(spPr, qn('a:ln'))
    tailEnd = ln.find(qn('a:tailEnd'))
    headEnd = ln.find(qn('a:headEnd'))
    if headEnd is None:
        headEnd = etree.SubElement(ln, qn('a:headEnd'))
    headEnd.set('type', 'none')
    if tailEnd is None:
        tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
    tailEnd.set('type', 'arrow')
    tailEnd.set('w', 'med')
    tailEnd.set('len', 'med')


def add_node(slide, cx, cy, rw, rh, label,
             fill=C_NODE_GREEN, text_color=C_WHITE, fsize=Pt(8.5)):
    """Centred ellipse node."""
    add_ellipse(slide,
                cx - rw/2, cy - rh/2, rw, rh,
                fill=fill, border=C_BLACK, border_w=Pt(0.8))
    add_textbox(slide,
                cx - rw/2, cy - rh/2, rw, rh,
                label, font_size=fsize, bold=True,
                color=text_color, align=PP_ALIGN.CENTER, wrap=True)


def add_entity_box(slide, x, y, w, h, label, fsize=Pt(8)):
    """Light-blue entity rectangle."""
    add_rect(slide, x, y, w, h,
             fill=C_BOX_BLUE_FILL, border=C_BOX_BLUE_BDR, border_w=Pt(1.0),
             radius=12000)
    add_textbox(slide, x, y, w, h,
                label, font_size=fsize, bold=False,
                color=C_BLACK, align=PP_ALIGN.CENTER, wrap=True)


def add_progress_bar(slide, x, y, w, h, fraction, bar_color=C_BAR_BLUE):
    """Background bar + filled portion."""
    add_rect(slide, x, y, w, h, fill=C_LIGHT_GRAY, border=None)
    if fraction > 0:
        add_rect(slide, x, y, int(w * fraction), h,
                 fill=bar_color, border=None)


# ═══════════════════════════════════════════════════════════════════════════
# BUILD SLIDE
# ═══════════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
slide = new_slide(prs)

# ── background ──────────────────────────────────────────────────────────────
bg = slide.shapes.add_shape(1, 0, 0, W, H)
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
bg.line.fill.background()

# ── title row (thin strip at very top) ──────────────────────────────────────
add_textbox(slide, Inches(0.1), Inches(0.02), W - Inches(0.2), Inches(0.3),
            "ProvMind: Graph-Grounded Process Memory for Materials Synthesis QA",
            font_size=Pt(13), bold=True,
            color=C_HDR_TEAL, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
#  PANEL GRID  (row0 y=0.32, row1 y=6.1  |  col0,1,2)
# ═══════════════════════════════════════════════════════════════════════════
TOP    = Inches(0.32)
MID    = Inches(6.05)
BOTTOM = H - Inches(0.08)
GAP    = Inches(0.12)

# column boundaries (x_start, width)
cols = [
    (Inches(0.08), Inches(6.50)),   # col 0
    (Inches(6.70), Inches(6.80)),   # col 1
    (Inches(13.62), Inches(6.30)),  # col 2
]

row_h_top = MID - TOP - GAP
row_h_bot = BOTTOM - MID - GAP

# ┌─────────────────────────────────────────────────────────────────────────┐
# │  PANEL A  – Graph-grounded query example                                │
# └─────────────────────────────────────────────────────────────────────────┘
ax, aw = cols[0]
ay, ah = TOP, row_h_top
add_label_box(slide, ax, ay, aw, ah,
              "A", "Graph-grounded query example",
              hdr_color=C_HDR_BLUE)

# --- graph nodes ---
# Row1: SiO2 substrate, Cr
ex, ey = ax + Inches(0.15), ay + Inches(0.55)
EW, EH = Inches(1.30), Inches(0.42)
NW, NH = Inches(1.35), Inches(0.46)  # activity node

add_entity_box(slide, ex,              ey, EW, EH, "SiO₂\nsubstrate")
add_entity_box(slide, ex + Inches(1.6),ey, EW, EH, "Cr")

# depositing (centre)
dcx = ax + Inches(2.65)
dcy = ay + Inches(1.25)
add_node(slide, dcx, dcy, NW, NH, "depositing")

# deposited sample
add_entity_box(slide, ax + Inches(1.80), ay + Inches(1.80), EW+Inches(0.1), EH, "deposited\nsample")

# Permalloy81, Aluminium oxide (left column)
add_entity_box(slide, ax + Inches(0.10), ay + Inches(1.60), EW, EH, "Permalloy81")
add_entity_box(slide, ax + Inches(0.10), ay + Inches(2.50), EW+Inches(0.1), EH, "Aluminum\noxide")

# sputtering node
scx = ax + Inches(2.65)
scy = ay + Inches(2.80)
add_node(slide, scx, scy, NW, NH, "sputtering")

# sputtered sample
add_entity_box(slide, ax + Inches(1.80), ay + Inches(3.45), EW+Inches(0.1), EH, "sputtered\nsample")

# second depositing
dcx2 = ax + Inches(2.65)
dcy2 = ay + Inches(4.35)
add_node(slide, dcx2, dcy2, NW, NH, "depositing")

# deposited sample 2
add_entity_box(slide, ax + Inches(1.80), ay + Inches(5.00), EW+Inches(0.1), EH, "deposited\nsample")

# Edwards diamond (tool)  – use rotated rect as diamond approximation
tool_cx = ax + Inches(4.70)
tool_cy = ay + Inches(1.00)
tool_w  = Inches(1.50)
tool_h  = Inches(0.60)
dmd = slide.shapes.add_shape(4, tool_cx - tool_w/2, tool_cy - tool_h/2,
                               tool_w, tool_h)  # 4=diamond
dmd.fill.solid()
dmd.fill.fore_color.rgb = C_DIAMOND_ORG
dmd.line.color.rgb = RGBColor(0xC0, 0x50, 0x00)
dmd.line.width = Pt(1.0)
add_textbox(slide, tool_cx - tool_w/2, tool_cy - tool_h/2 - Inches(0.05),
            tool_w, tool_h + Inches(0.1),
            "Edwards\n2-inch\nmagnetron\nsputtering\nsystem",
            font_size=Pt(6.5), bold=False,
            color=C_BLACK, align=PP_ALIGN.CENTER)

# condition boxes (grey)
def cond_box(label, x, y):
    add_rect(slide, x, y, Inches(1.55), Inches(0.52),
             fill=C_LIGHT_GRAY, border=RGBColor(0x80,0x80,0x80), border_w=Pt(0.8))
    add_textbox(slide, x, y, Inches(1.55), Inches(0.52),
                label, font_size=Pt(7.5), color=C_BLACK, align=PP_ALIGN.LEFT)

cond_box("length_thickness: 50 nm\natmosphere: Ar\npressure: 4 mTorr",
         ax + Inches(3.65), ay + Inches(1.55))
cond_box("concentration: 81% Ni/\n17% Fe/2% Mo",
         ax + Inches(0.10), ay + Inches(2.20))
cond_box("atmosphere: Ar\npressure: 4 mTorr",
         ax + Inches(3.65), ay + Inches(2.90))
cond_box("atmosphere: Ar\npressure: 4 mTorr",
         ax + Inches(3.65), ay + Inches(4.55))

# --- question + options ---
qy = ay + ah - Inches(1.70)
add_textbox(slide, ax + Inches(0.10), qy, aw - Inches(0.20), Inches(0.50),
            "Question: For step 2, sputtering, what is the required atmosphere / gas environment?",
            font_size=Pt(9), bold=False, italic=True,
            color=C_BLACK, align=PP_ALIGN.LEFT)

opts = [("A", "Vacuum", False), ("B", "N₂", False),
        ("C", "Ar", True),      ("D", "Ar + O₂", False)]
ox_start = [ax + Inches(0.20), ax + Inches(3.30),
            ax + Inches(0.20), ax + Inches(3.30)]
oy_start = [qy + Inches(0.60)] * 2 + [qy + Inches(1.15)] * 2
for i, (ltr, txt, correct) in enumerate(opts):
    ox, oy = ox_start[i], oy_start[i]
    ow, oh = Inches(2.80), Inches(0.44)
    fc = C_GREEN_TICK if correct else C_WHITE
    bc = C_GREEN_TICK if correct else C_BOX_BLUE_BDR
    add_rect(slide, ox, oy, ow, oh,
             fill=fc, border=bc, border_w=Pt(1.5), radius=15000)
    label = f"{ltr}   {txt}" + ("  ✓" if correct else "")
    add_textbox(slide, ox + Inches(0.08), oy, ow - Inches(0.1), oh,
                label, font_size=Pt(9.5), bold=correct,
                color=C_WHITE if correct else C_BLACK,
                align=PP_ALIGN.LEFT)

# ┌─────────────────────────────────────────────────────────────────────────┐
# │  PANEL B  – Stage 0: Offline process memory construction                │
# └─────────────────────────────────────────────────────────────────────────┘
bx, bw = cols[1]
by, bh = TOP, row_h_top
add_label_box(slide, bx, by, bw, bh,
              "B", "Stage 0: Offline process memory construction",
              hdr_color=C_HDR_TEAL)

# Raw JSON-LD box
rjx = bx + Inches(0.20)
rjy = by + Inches(0.55)
rjw, rjh = Inches(1.50), Inches(1.00)
add_rect(slide, rjx, rjy, rjw, rjh,
         fill=C_BOX_BLUE_FILL, border=C_BOX_BLUE_BDR, border_w=Pt(1.0), radius=10000)
add_textbox(slide, rjx, rjy, rjw, rjh,
            "Raw JSON-LD\nprovenance\ngraphs",
            font_size=Pt(8.5), bold=False,
            color=C_BLACK, align=PP_ALIGN.CENTER)
# JSON icon
add_rect(slide, rjx + Inches(0.15), rjy + Inches(0.05),
         Inches(0.55), Inches(0.40),
         fill=RGBColor(0x17,0x4A,0x7E), border=None)
add_textbox(slide, rjx + Inches(0.15), rjy + Inches(0.05),
            Inches(0.55), Inches(0.40),
            "{...}", font_size=Pt(9), bold=True,
            color=C_WHITE, align=PP_ALIGN.CENTER)

# compiler.py
add_textbox(slide, rjx + Inches(1.65), rjy + Inches(0.15),
            Inches(1.10), Inches(0.40),
            "compiler.py",
            font_size=Pt(8.5), bold=False, italic=True,
            color=C_BLACK, align=PP_ALIGN.CENTER)
add_arrow(slide,
          rjx + rjw, rjy + rjh/2,
          rjx + rjw + Inches(1.00), rjy + rjh/2)

# ProcessState box
psx = bx + Inches(2.85)
psy = rjy
psw, psh = Inches(2.20), Inches(1.60)
add_rect(slide, psx, psy, psw, psh,
         fill=C_BOX_BLUE_FILL, border=C_BOX_BLUE_BDR, border_w=Pt(1.0), radius=10000)
add_textbox(slide, psx + Inches(0.05), psy + Inches(0.02), psw - Inches(0.1), Inches(0.35),
            "ProcessState", font_size=Pt(9), bold=True,
            color=C_BLACK, align=PP_ALIGN.LEFT)
for i, line in enumerate(["• entities", "• activities", "• usage", "• generation"]):
    add_textbox(slide, psx + Inches(0.12), psy + Inches(0.38) + i*Inches(0.27),
                psw - Inches(0.15), Inches(0.28),
                line, font_size=Pt(8.5),
                color=C_BLACK, align=PP_ALIGN.LEFT)

# arrow to ProcessMemoryIndex
add_arrow(slide, psx + psw, psy + psh/2,
          psx + psw + Inches(0.35), psy + psh/2)

# ProcessMemoryIndex (database cylinder styled as rect)
pmx = bx + Inches(5.35)
pmy = rjy
pmw, pmh = Inches(1.20), Inches(1.60)
add_rect(slide, pmx, pmy, pmw, pmh,
         fill=RGBColor(0xD5, 0xBC, 0xF0), border=C_DB_PURPLE, border_w=Pt(1.2), radius=15000)
add_textbox(slide, pmx, pmy, pmw, pmh,
            "ProcessMemoryIndex\n(train only)",
            font_size=Pt(8), bold=True,
            color=C_BLACK, align=PP_ALIGN.CENTER)

# Kahn / role labels
add_textbox(slide, bx + Inches(0.10), by + Inches(2.35), bw - Inches(0.2), Inches(0.30),
            "Kahn topological sort → order",
            font_size=Pt(8), italic=True,
            color=RGBColor(0x40,0x40,0x40), align=PP_ALIGN.LEFT)
add_textbox(slide, bx + Inches(0.10), by + Inches(2.65), bw - Inches(0.2), Inches(0.30),
            "role assignment → precursor / intermediate / product",
            font_size=Pt(8), italic=True,
            color=RGBColor(0x40,0x40,0x40), align=PP_ALIGN.LEFT)

# Three index boxes at bottom
idx_labels = [("global_bigrams\n(counts)", C_BOX_BLUE_FILL),
              ("prefix_next\n(transitions)", C_BOX_BLUE_FILL),
              ("step_library\n(metadata)", C_BOX_BLUE_FILL)]
for i, (lbl, fc) in enumerate(idx_labels):
    ix = bx + Inches(0.15) + i * Inches(2.20)
    iy = by + bh - Inches(1.80)
    add_rect(slide, ix, iy, Inches(2.00), Inches(0.80),
             fill=fc, border=C_BOX_BLUE_BDR, border_w=Pt(1.0), radius=10000)
    add_textbox(slide, ix, iy, Inches(2.00), Inches(0.80),
                lbl, font_size=Pt(8.5),
                color=C_BLACK, align=PP_ALIGN.CENTER)
    # arrow from ProcessMemoryIndex down
    add_arrow(slide, pmx + pmw/2, pmy + pmh,
              ix + Inches(1.00), iy,
              color=RGBColor(0x50,0x50,0x50))

# ┌─────────────────────────────────────────────────────────────────────────┐
# │  PANEL C  – Stage 1: Analogy retrieval                                  │
# └─────────────────────────────────────────────────────────────────────────┘
cx2, cw = cols[2]
cy2, ch = TOP, row_h_top
add_label_box(slide, cx2, cy2, cw, ch,
              "C", "Stage 1: Analogy retrieval",
              hdr_color=C_HDR_TEAL)

# Query process (left column)
add_textbox(slide, cx2 + Inches(0.10), cy2 + Inches(0.45),
            Inches(1.20), Inches(0.30),
            "Query process", font_size=Pt(8.5), bold=True,
            color=C_BLACK, align=PP_ALIGN.CENTER)
qproc_nodes = ["depositing", "sputtering", "depositing"]
for i, n in enumerate(qproc_nodes):
    add_node(slide,
             cx2 + Inches(0.65),
             cy2 + Inches(0.90) + i * Inches(0.95),
             Inches(1.20), Inches(0.42), n,
             fsize=Pt(8))
    if i < 2:
        add_arrow(slide,
                  cx2 + Inches(0.65),
                  cy2 + Inches(1.12) + i * Inches(0.95),
                  cx2 + Inches(0.65),
                  cy2 + Inches(1.12) + (i+1) * Inches(0.95) - Inches(0.20))

# Heuristic retrieval score box
hx = cx2 + Inches(1.55)
hy = cy2 + Inches(0.45)
hw, hh = Inches(2.70), Inches(1.45)
add_rect(slide, hx, hy, hw, hh,
         fill=C_BOX_BLUE_FILL, border=C_BOX_BLUE_BDR, border_w=Pt(1.0), radius=10000)
add_textbox(slide, hx + Inches(0.05), hy + Inches(0.04), hw - Inches(0.1), Inches(0.30),
            "Heuristic retrieval score", font_size=Pt(8.5), bold=True,
            color=C_BLACK, align=PP_ALIGN.LEFT)
score_formula = ("score = 3 × |query steps ∩ route steps|\n"
                 "      + 4 × 1[length match]\n"
                 "      + |material token overlap|")
add_textbox(slide, hx + Inches(0.08), hy + Inches(0.36), hw - Inches(0.15), Inches(1.00),
            score_formula, font_size=Pt(8),
            color=C_BLACK, align=PP_ALIGN.LEFT)

# Top-k retrieved processes box
rx = hx
ry = hy + hh + Inches(0.15)
rw2, rh2 = hw, Inches(2.00)
add_rect(slide, rx, ry, rw2, rh2,
         fill=C_BOX_BLUE_FILL, border=C_BOX_BLUE_BDR, border_w=Pt(1.0), radius=10000)
add_textbox(slide, rx + Inches(0.05), ry + Inches(0.04), rw2 - Inches(0.1), Inches(0.30),
            "Top-k retrieved training processes", font_size=Pt(8.5), bold=True,
            color=C_BLACK, align=PP_ALIGN.LEFT)
retrieved = [
    "1  depositing → sputtering → depositing",
    "2  preparing → sputtering → annealing",
    "3  deposition → sputtering → deposition",
]
for i, line in enumerate(retrieved):
    add_textbox(slide, rx + Inches(0.08), ry + Inches(0.38) + i * Inches(0.48),
                rw2 - Inches(0.15), Inches(0.40),
                line, font_size=Pt(8.5),
                color=C_BLACK, align=PP_ALIGN.LEFT)

# Train processes (right column of C)
add_textbox(slide, cx2 + Inches(4.75), cy2 + Inches(0.45),
            Inches(1.40), Inches(0.30),
            "Train processes\n(routes)", font_size=Pt(8), bold=True,
            color=C_BLACK, align=PP_ALIGN.CENTER)
# small chain icons
for row in range(3):
    for col in range(3):
        nx = cx2 + Inches(4.80) + col * Inches(0.42)
        ny = cy2 + Inches(0.82) + row * Inches(0.65)
        fc = C_NODE_GREEN if col == 1 else C_BOX_BLUE_FILL
        add_ellipse(slide, nx, ny, Inches(0.30), Inches(0.24),
                    fill=fc, border=C_BLACK, border_w=Pt(0.5))
        if col < 2:
            add_arrow(slide, nx + Inches(0.30), ny + Inches(0.12),
                      nx + Inches(0.42), ny + Inches(0.12),
                      width=Pt(0.8))

# DualView extension box
dvx = cx2 + Inches(0.10)
dvy = cy2 + ch - Inches(1.60)
dvw = cw - Inches(0.20)
dvh = Inches(1.45)
add_rect(slide, dvx, dvy, dvw, dvh,
         fill=RGBColor(0xEF, 0xF5, 0xFF), border=C_HDR_BLUE, border_w=Pt(1.5),
         radius=10000)
add_textbox(slide, dvx + Inches(0.08), dvy + Inches(0.04), dvw - Inches(0.15), Inches(0.30),
            "DualView extension (optional)", font_size=Pt(8.5), bold=True,
            color=C_HDR_BLUE, align=PP_ALIGN.LEFT)
add_textbox(slide, dvx + Inches(0.08), dvy + Inches(0.36), dvw - Inches(0.15), Inches(0.40),
            "Fused retrieval score:  α·text_sim + β·struct_sim + γ·heuristic_score",
            font_size=Pt(8.5), italic=False,
            color=C_BLACK, align=PP_ALIGN.LEFT)
add_textbox(slide, dvx + Inches(0.20), dvy + Inches(0.82), Inches(2.60), Inches(0.38),
            "📄  Text view: SentenceTransformer",
            font_size=Pt(8.5),
            color=C_BLACK, align=PP_ALIGN.LEFT)
add_textbox(slide, dvx + Inches(3.30), dvy + Inches(0.82), Inches(2.70), Inches(0.38),
            "🕸  Structure view: frozen 2-layer GAT",
            font_size=Pt(8.5),
            color=C_BLACK, align=PP_ALIGN.LEFT)

# ┌─────────────────────────────────────────────────────────────────────────┐
# │  PANEL D  – Stage 2: Task-aware symbolic scorer                         │
# └─────────────────────────────────────────────────────────────────────────┘
dx, dw = cols[0]
dy, dh = MID, row_h_bot
add_label_box(slide, dx, dy, dw, dh,
              "D", "Stage 2: Task-aware symbolic scorer",
              hdr_color=C_HDR_TEAL)

# Section rows: A Route, B Condition, C Tool, D Order
sections = [
    ("A", "Route-level reasoning",    C_SECTION_A),
    ("B", "Condition reasoning",      C_SECTION_B),
    ("C", "Tool reasoning",           C_SECTION_C),
    ("D", "Order reasoning",          C_SECTION_D2),
]
sub_items = {
    "A": [("A1", "Route retrieval",   "bigrams + sequence similarity"),
          ("A2", "Missing step",      "position / form / condition voting"),
          ("A3", "Next activity",     "prefix_next + bigrams + retrieved routes")],
    "B": [("B1", "Single condition",  "matched-step condition voting"),
          ("B2", "Full condition set","(temperature, duration, atmosphere)\ntriplet voting")],
    "C": [("C1", "Tool selection",    "tool frequency voting")],
    "D": [("D1", "Process ordering",  "topological constraint satisfaction")],
}

sy = dy + Inches(0.45)
for sec_ltr, sec_title, sec_color in sections:
    # section header bar
    add_rect(slide, dx + Inches(0.12), sy, dw - Inches(0.24), Inches(0.30),
             fill=sec_color, border=None, radius=8000)
    add_textbox(slide, dx + Inches(0.18), sy, dw - Inches(0.30), Inches(0.30),
                f"{sec_ltr}   {sec_title}",
                font_size=Pt(8.5), bold=True,
                color=C_WHITE if sec_ltr != "C" else C_BLACK,
                align=PP_ALIGN.LEFT)
    sy += Inches(0.30)
    for code, name, desc in sub_items[sec_ltr]:
        add_textbox(slide, dx + Inches(0.22), sy, Inches(1.20), Inches(0.30),
                    code, font_size=Pt(8.5), bold=True,
                    color=C_BLACK, align=PP_ALIGN.LEFT)
        add_textbox(slide, dx + Inches(0.60), sy, Inches(1.80), Inches(0.30),
                    name, font_size=Pt(8.5),
                    color=C_BLACK, align=PP_ALIGN.LEFT)
        add_textbox(slide, dx + Inches(2.55), sy, Inches(3.80), Inches(0.30),
                    desc, font_size=Pt(8), italic=True,
                    color=RGBColor(0x40,0x40,0x40), align=PP_ALIGN.LEFT)
        sy += Inches(0.30)
    sy += Inches(0.04)

# Example option scores table
ty = dy + dh - Inches(2.00)
add_textbox(slide, dx + Inches(0.12), ty, dw - Inches(0.24), Inches(0.28),
            "Example option scores (step 2 atmosphere)",
            font_size=Pt(8.5), bold=True,
            color=C_BLACK, align=PP_ALIGN.LEFT)
ty += Inches(0.28)
# header row
for col_txt, col_x, col_w in [("Option", dx+Inches(0.15), Inches(0.60)),
                                ("", dx+Inches(0.78), Inches(3.80)),
                                ("Symbolic score (0–1)", dx+Inches(0.78), Inches(3.80))]:
    pass
add_textbox(slide, dx+Inches(0.15), ty, Inches(0.60), Inches(0.25),
            "Option", font_size=Pt(8), bold=True,
            color=C_BLACK, align=PP_ALIGN.LEFT)
add_textbox(slide, dx+Inches(0.78), ty, Inches(3.80), Inches(0.25),
            "Symbolic score (0–1)", font_size=Pt(8), bold=True,
            color=C_BLACK, align=PP_ALIGN.LEFT)
ty += Inches(0.26)

score_rows = [("A", "Vacuum", 0.12, False),
              ("B", "N₂",     0.08, False),
              ("C", "Ar",     0.94, True),
              ("D", "Ar + O₂",0.15, False)]
for opt, name, score, correct in score_rows:
    fc = RGBColor(0xE2, 0xEF, 0xDA) if correct else C_WHITE
    add_rect(slide, dx + Inches(0.12), ty, dw - Inches(0.24), Inches(0.30),
             fill=fc, border=C_LIGHT_GRAY, border_w=Pt(0.5))
    add_textbox(slide, dx+Inches(0.15), ty, Inches(0.25), Inches(0.30),
                opt, font_size=Pt(8.5), bold=correct,
                color=C_BLACK, align=PP_ALIGN.CENTER)
    add_textbox(slide, dx+Inches(0.45), ty, Inches(1.10), Inches(0.30),
                name + (" ✓" if correct else ""),
                font_size=Pt(8.5), bold=correct,
                color=C_GREEN_TICK if correct else C_BLACK,
                align=PP_ALIGN.LEFT)
    bar_x = dx + Inches(1.65)
    bar_y = ty + Inches(0.05)
    add_progress_bar(slide, bar_x, bar_y, Inches(3.50), Inches(0.20), score,
                     bar_color=C_GREEN_TICK if correct else C_BAR_BLUE)
    add_textbox(slide, bar_x + Inches(3.55), ty, Inches(0.45), Inches(0.28),
                f"{score:.2f}", font_size=Pt(8),
                color=C_BLACK, align=PP_ALIGN.LEFT)
    ty += Inches(0.30)

# scale: 0 … 0.5 … 1.0
add_textbox(slide, dx+Inches(1.65), ty, Inches(0.30), Inches(0.22),
            "0", font_size=Pt(7.5), color=RGBColor(0x60,0x60,0x60), align=PP_ALIGN.LEFT)
add_textbox(slide, dx+Inches(3.35), ty, Inches(0.50), Inches(0.22),
            "0.5", font_size=Pt(7.5), color=RGBColor(0x60,0x60,0x60), align=PP_ALIGN.CENTER)
add_textbox(slide, dx+Inches(5.10), ty, Inches(0.40), Inches(0.22),
            "1.0", font_size=Pt(7.5), color=RGBColor(0x60,0x60,0x60), align=PP_ALIGN.RIGHT)

# ┌─────────────────────────────────────────────────────────────────────────┐
# │  PANEL E  – Stages 3-5: LLM planning, decision, and fallback            │
# └─────────────────────────────────────────────────────────────────────────┘
ex2, ew = cols[1]
ey2, eh = MID, row_h_bot
add_label_box(slide, ex2, ey2, ew, eh,
              "E", "Stages 3–5: LLM planning, decision, and fallback",
              hdr_color=C_HDR_TEAL)

stages = [
    ("1", "LLM Planning",  C_HDR_BLUE,
     "Prompt contains:\n• question\n• top-5 retrieved process summaries\n• symbolic option scores\n"
     "Model: Qwen2.5-7B-Instruct\nOutput: 2–4 step plan (max 96 tokens)"),
    ("2", "LLM Decision",  C_HDR_TEAL,
     "Given plan + evidence,\nchoose the best option letter.\n"
     "Output: Answer: X\n(max 48 tokens)"),
    ("3", "Symbolic\nfallback", C_SECTION_B,
     "If model output is invalid:\n"
     "  if valid answer letter → use it\n"
     "  else → use symbolic_best"),
]
svy = ey2 + Inches(0.48)
for num, title, color, body in stages:
    # numbered circle
    circ = slide.shapes.add_shape(9,
                                   ex2 + Inches(0.15), svy,
                                   Inches(0.36), Inches(0.36))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    add_textbox(slide, ex2 + Inches(0.15), svy,
                Inches(0.36), Inches(0.36),
                num, font_size=Pt(10), bold=True,
                color=C_WHITE, align=PP_ALIGN.CENTER)
    # title bar
    add_rect(slide, ex2 + Inches(0.58), svy + Inches(0.02),
             ew - Inches(0.78), Inches(0.32),
             fill=color, border=None, radius=8000)
    add_textbox(slide, ex2 + Inches(0.65), svy + Inches(0.02),
                ew - Inches(0.85), Inches(0.32),
                title, font_size=Pt(9.5), bold=True,
                color=C_WHITE, align=PP_ALIGN.LEFT)
    # body
    add_textbox(slide, ex2 + Inches(0.18), svy + Inches(0.38),
                ew - Inches(0.30), Inches(1.15),
                body, font_size=Pt(8.5),
                color=C_BLACK, align=PP_ALIGN.LEFT)
    svy += Inches(1.60)

# ┌─────────────────────────────────────────────────────────────────────────┐
# │  PANEL F  – ProvMind output                                             │
# └─────────────────────────────────────────────────────────────────────────┘
fx, fw = cols[2]
fy, fh = MID, row_h_bot
add_label_box(slide, fx, fy, fw, fh,
              "F", "ProvMind output",
              hdr_color=C_HDR_TEAL)

# "Predicted condition value:" label
add_textbox(slide, fx + Inches(0.20), fy + Inches(0.52),
            fw - Inches(0.40), Inches(0.35),
            "Predicted condition value:",
            font_size=Pt(10.5), bold=True,
            color=C_HDR_TEAL, align=PP_ALIGN.CENTER)

# big "Ar" answer
ar_box = add_rect(slide, fx + fw/2 - Inches(1.00), fy + Inches(1.00),
                  Inches(2.00), Inches(1.10),
                  fill=RGBColor(0xE8, 0xF5, 0xE9),
                  border=C_GREEN_TICK, border_w=Pt(2), radius=15000)
add_textbox(slide, fx + fw/2 - Inches(1.00), fy + Inches(1.00),
            Inches(2.00), Inches(1.10),
            "Ar", font_size=Pt(48), bold=True,
            color=C_GREEN_TICK, align=PP_ALIGN.CENTER)

# tick-list
checks = [
    "retrieved analogs agree on Ar",
    "symbolic option score highest for Ar",
    "LLM plan grounded in\ntrain-process evidence",
]
cky = fy + Inches(2.25)
for ck in checks:
    add_textbox(slide, fx + Inches(0.25), cky,
                fw - Inches(0.50), Inches(0.52),
                "✅  " + ck, font_size=Pt(9),
                color=C_BLACK, align=PP_ALIGN.LEFT)
    cky += Inches(0.52)

# train-memory label at bottom
tmx = fx + Inches(0.30)
tmy = fy + fh - Inches(0.90)
add_rect(slide, tmx, tmy, fw - Inches(0.60), Inches(0.72),
         fill=RGBColor(0xD5, 0xF5, 0xD5),
         border=C_GREEN_TICK, border_w=Pt(1.5), radius=12000)
add_textbox(slide, tmx, tmy, fw - Inches(0.60), Inches(0.72),
            "🔒  train-memory only\n(fair train-process only)",
            font_size=Pt(9), bold=False,
            color=C_GREEN_TICK, align=PP_ALIGN.CENTER)

# ─── save ───────────────────────────────────────────────────────────────────
outpath = "provmind_figure.pptx"
prs.save(outpath)
print(f"Saved → {outpath}")
